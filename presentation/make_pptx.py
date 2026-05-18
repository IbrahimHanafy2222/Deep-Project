"""
Generate the project presentation as an editable PowerPoint deck.

Usage:
    pip install python-pptx
    python make_pptx.py        # writes slides.pptx next to this script

Content mirrors slides.tex / slides.html. Figures are pulled from
../report/figures and ../evaluation/phase5_transformer.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

HERE = os.path.dirname(os.path.abspath(__file__))
P5DIR = os.path.join(HERE, "..", "evaluation", "phase5_transformer")
FIG_BAHDANAU = os.path.join(HERE, "..", "report", "figures", "attn_3.png")
FIG_XATTN = os.path.join(P5DIR, "enc_self_attn_3.png")
FIG_CROSS = os.path.join(P5DIR, "dec_cross_attn_3.png")
FIG_LEN = os.path.join(P5DIR, "three_way_compare.png")

NAVY = RGBColor(0x1D, 0x3B, 0x6E)
RUST = RGBColor(0xB8, 0x35, 0x0F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def title_bar(slide, text):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35),
                                   Inches(12.1), Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = NAVY


def bullets(slide, items, top=1.5, left=0.8, width=11.7, size=22):
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(width), Inches(5.0))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, level) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = Pt(size if level == 0 else size - 4)
        p.space_after = Pt(8)


def box_note(slide, text, top=5.2, left=0.8, width=11.7, size=20):
    shape = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(1.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xEE, 0xF2, 0xFB)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = NAVY


# ---------------------------------------------------------------- 1 title
s = add_slide()
t = s.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(2.0))
p = t.text_frame.paragraphs[0]
p.text = "Grammatical Error Correction on C4_200M"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = NAVY
sub = s.shapes.add_textbox(Inches(0.8), Inches(3.6), Inches(11.7), Inches(2.5))
tf = sub.text_frame
tf.word_wrap = True
for i, line in enumerate([
        "Bag-of-Words  ->  LSTM  ->  Attention  ->  Transformer", "",
        "Abdullah Ahmed (202200206)   |   Ibrahim Hanafy (202200518)",
        "CIE 555 - Neural Networks and Deep Learning",
        "Instructor: Dr. Ibrahim Swelam   |   TAs: Aya Abdelaziz, Mahmoud Farahat"]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    p.font.size = Pt(20)

# ---------------------------------------------------------------- 2 problem
s = add_slide()
title_bar(s, "The Problem: Grammatical Error Correction")
bullets(s, [
    ("Rewrite an ungrammatical sentence into its corrected form.", 0),
    ("Errors: tense, agreement, articles, plurals, punctuation, spelling.", 0),
    ("The output is a whole new sentence - a sequence-to-sequence task.", 0),
])
box_note(s, 'Input:  "She go to school yesterday and learning many thing."\n'
            'Output: "She went to school yesterday and learned many things."',
         top=4.6)

# ---------------------------------------------------------------- 3 data
s = add_slide()
title_bar(s, "The Data: C4_200M Subset")
bullets(s, [
    ("Clean web text automatically corrupted into (corrupted, clean) pairs.", 0),
    ("Train: 850,000   Validation: 50,000   Test: 100,000   (Total 1,000,000)", 0),
    ("16k Byte-Level BPE tokenizer; max length 64 tokens.", 0),
    ("Same splits and tokenizer for every model - no data leakage.", 0),
])

# ---------------------------------------------------------------- 4 four models
s = add_slide()
title_bar(s, "Four Models of Increasing Power")
rows = [("#", "Model", "What it adds"),
        ("1", "Bag-of-Words (TF-IDF)", "Baseline - only detects errors"),
        ("2", "LSTM encoder-decoder", "First corrector; fixed-size bottleneck"),
        ("3", "LSTM + Bahdanau attention", "Decoder re-queries the full source"),
        ("4", "Transformer", "Self-attention; no recurrence")]
tbl = s.shapes.add_table(5, 3, Inches(1.0), Inches(1.7),
                         Inches(11.3), Inches(3.2)).table
tbl.columns[0].width = Inches(0.9)
tbl.columns[1].width = Inches(4.4)
tbl.columns[2].width = Inches(6.0)
for r, row in enumerate(rows):
    for c, val in enumerate(row):
        cell = tbl.cell(r, c)
        cell.text = val
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(18)
        if r == 0:
            para.font.bold = True
            para.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
box_note(s, "Model 1 detects.  Models 2-4 correct.", top=5.3, size=20)

# ---------------------------------------------------------------- 5 BoW
s = add_slide()
title_bar(s, "Model 1 - Bag-of-Words Baseline")
bullets(s, [
    ("TF-IDF (unigrams+bigrams, 50k vocab) + Logistic Regression / Linear SVM. "
     "Each pair -> corrupted = 0, clean = 1.", 0),
    ("Result: ~62% accuracy - barely above 50% chance. Why so weak? "
     "Structural, not bad tuning:", 0),
    ("Detector, not corrector - outputs one label, never corrected text.", 1),
    ("Bag-of-words discards order - but grammar IS order "
     '("she goes" vs "go she").', 1),
    ("Error signal swamped - a 1-word error is invisible next to "
     "content-word weights.", 1),
    ("No generalisation - only memorises lexical statistics.", 1),
])
box_note(s, "Proves GEC needs sequence modelling and text generation - "
            "which the next three models provide.", top=5.3, size=18)

# ---------------------------------------------------------------- 6 LSTM
s = add_slide()
title_bar(s, "Model 2 - LSTM Encoder-Decoder")
bullets(s, [
    ("BiLSTM encoder compresses the input into ONE 512-dim state.", 0),
    ("LSTM decoder generates the correction token by token.", 0),
    ("Teacher forcing 1.0 -> 0.5; 19.55M parameters; 3 epochs.", 0),
    ("Test GLEU 0.213 (greedy) - the weakest corrector.", 0),
])
box_note(s, "The bottleneck: after the first step the decoder never sees the "
            "source again - everything is squeezed into 512 numbers.", top=4.7)

# ---------------------------------------------------------------- 7 Bahdanau
s = add_slide()
title_bar(s, "Model 3 - LSTM + Bahdanau Attention")
bullets(s, [
    ("Keep ALL encoder states; the decoder attends to them every step.", 0),
    ("score(s,h) = v.tanh(W.s + W.h)  ->  alpha = softmax  ->  "
     "context = sum(alpha_i . h_i).", 0),
    ("Removes the bottleneck; alignment weights are interpretable.", 0),
    ("Padding masked with -inf before softmax. 29.06M parameters.", 0),
])
box_note(s, "The decisive jump: GLEU 0.213 -> 0.599 (+0.386). "
            "Attention matters most.", top=4.7)

# ---------------------------------------------------------------- 8 Transformer
s = add_slide()
title_bar(s, "Model 4 - Transformer")
bullets(s, [
    ("No recurrence - self-attention only.", 0),
    ("d_model 256, 8 heads, 3+3 layers, FFN 512.", 0),
    ("Sinusoidal positional encoding; causal mask in the decoder.", 0),
    ("Label smoothing 0.1; warmup learning rate; weight tying.", 0),
])
box_note(s, "Smallest and strongest: only 8.05M parameters (3.6x fewer than "
            "Model 3). Best validation GLEU 0.613.", top=4.7)

# ---------------------------------------------------------------- 9 results
s = add_slide()
title_bar(s, "Comparative Results - Test Set")
rows = [("Model", "Decoding", "Params", "EM", "GLEU"),
        ("LSTM (no attention)", "greedy", "19.55M", "~0.008", "0.213"),
        ("LSTM (no attention)", "beam-4", "19.55M", "-", "0.256"),
        ("LSTM + Bahdanau", "greedy", "29.06M", "0.024", "0.599"),
        ("LSTM + Bahdanau", "beam-4", "29.06M", "0.026", "0.618"),
        ("Transformer", "greedy", "8.05M", "0.022", "0.618"),
        ("Transformer", "beam-4", "8.05M", "0.020", "0.618")]
tbl = s.shapes.add_table(7, 5, Inches(1.6), Inches(1.6),
                         Inches(10.1), Inches(3.6)).table
for r, row in enumerate(rows):
    for c, val in enumerate(row):
        cell = tbl.cell(r, c)
        cell.text = val
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(16)
        if r == 0:
            para.font.bold = True
            para.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
box_note(s, "Beam search helps the LSTM, not the well-calibrated Transformer. "
            "Transformer ties Model 3 on GLEU with 3.6x fewer parameters.",
         top=5.4, size=18)

# ---------------------------------------------------------------- 10 attention figs
s = add_slide()
title_bar(s, "Attention Is Interpretable")
for fig, left in [(FIG_BAHDANAU, 0.5), (FIG_XATTN, 4.7), (FIG_CROSS, 8.9)]:
    if os.path.exists(fig):
        s.shapes.add_picture(fig, Inches(left), Inches(1.5), height=Inches(3.7))
box_note(s, "Bahdanau decoder alignment | Transformer encoder self-attention | "
            "Transformer decoder cross-attention. Bahdanau and Transformer "
            "cross-attention learn visually similar source-to-output alignments.",
         top=5.5, size=17)

# ---------------------------------------------------------------- 11 length
s = add_slide()
title_bar(s, "Behaviour by Sentence Length")
if os.path.exists(FIG_LEN):
    s.shapes.add_picture(FIG_LEN, Inches(3.6), Inches(1.5), height=Inches(3.4))
bullets(s, [
    ("Bottleneck LSTM (blue) collapses to ~0 on medium/long sentences.", 0),
    ("Bahdanau and Transformer stay non-zero - attention prevents the collapse.", 0),
    ("Crossover: Transformer (green) overtakes Bahdanau (orange) on the LONG "
     "bucket - its O(1) path helps most there.", 0),
], top=5.1, size=18)

# ---------------------------------------------------------------- 11b worked example
s = add_slide()
title_bar(s, "Worked Example - The Three Correctors")
box = s.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12.0), Inches(5.4))
tf = box.text_frame
tf.word_wrap = True
lines = [
    ("Short sentence (insert a missing article):", True),
    ("SRC : Why new record is taking so long.", False),
    ("REF : Why the new record is taking so long.", False),
    ("P3  : Why new record is taking so long.    LSTM   - no edit (copies)", False),
    ("P4  : Why a new record is taking so long.  Bahd.  - inserts article 'a'", False),
    ("P5  : Why new record is taking so long?    Trans. - '.' -> '?'", False),
    ("", False),
    ("Long, noisy sentence:", True),
    ("P3 : '...cpue to buy and/indukine stead...'  hallucinated nonsense", False),
    ("P4 : copies source, corrupts the URL", False),
    ("P5 : copies source, fixes 'got' -> 'get'     the one real fix", False),
    ("", False),
    ("LSTM collapses on long input; P4 and P5 differ by only 1-2 tokens.", True),
]
for i, (text, bold) in enumerate(lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = text
    p.font.bold = bold
    p.font.size = Pt(19) if bold else Pt(16)
    if not bold and text:
        p.font.name = "Consolas"

# ---------------------------------------------------------------- 11c why P4 ~ P5
s = add_slide()
title_bar(s, "Why Bahdanau LSTM ~ Transformer")
bullets(s, [
    ("Test GLEU is tied at 0.618. The Transformer matches, not beats:", 0),
    ("Decisive feature is shared - both have full source access via attention "
     "(the +0.386 jump). LSTM-vs-self-attention is second-order.", 1),
    ("Both hit the data ceiling - same budget, same noisy C4_200M references "
     "-> both converge to GLEU ~0.62.", 1),
    ("Sentences are mostly short - the O(1) path only helps long-range "
     "dependencies (see the crossover).", 1),
    ("On hard inputs both just copy - outputs differ by 1-2 tokens.", 1),
])
box_note(s, "The real win is efficiency: same quality with 8.05M params vs "
            "29.06M (3.6x fewer) + faster parallel training - it would pull "
            "ahead with more data or cleaner references.", top=5.0, size=18)

# ---------------------------------------------------------------- 12 grammar
s = add_slide()
title_bar(s, "Grammar-Rule Analysis of Outputs")
bullets(s, [
    ("Learned well (local syntax):", 0),
    ('Subject-verb agreement: "What are this job" -> "What is this job"', 1),
    ('Articles: "Why new record" -> "Why the new record"', 1),
    ("Verb tense, punctuation, spacing.", 1),
    ("Still hard:", 0),
    ("Long-range / semantic edits; rare-word spelling (hallucination).", 1),
    ("Dominant failure = under-correction: copying the source unchanged.", 1),
])

# ---------------------------------------------------------------- 13 challenge
s = add_slide()
title_bar(s, "The Final Challenge - Dataset Quality")
bullets(s, [
    ("Why does EVERY model score only ~2% exact match?", 0),
    ("Reference quality is the ceiling - not model capacity:", 0),
    ('References are paraphrases: "informed"->"suggest", '
     '"startups"->"start-ups" - not grammar errors.', 1),
    ('References can be noisy: target "no brain worky" - not a word.', 1),
    ("Many valid corrections exist; exact match credits exactly one.", 1),
])
box_note(s, "Report GLEU / F0.5, not exact match. Cleaner data would help "
            "more than a bigger model.", top=5.3, size=19)

# ---------------------------------------------------------------- 14 conclusion
s = add_slide()
title_bar(s, "Conclusion")
bullets(s, [
    ("BoW detects grammaticality at ~62% - cannot correct.", 0),
    ("LSTM corrects but is limited by the bottleneck (GLEU 0.213).", 0),
    ("Bahdanau attention is the decisive jump (GLEU 0.599).", 0),
    ("Transformer ties on GLEU (0.618) with 3.6x fewer params; "
     "best validation GLEU 0.613.", 0),
])
box_note(s, "Key insight: the binding constraint is dataset quality, not "
            "architecture - synthetic C4_200M references are often "
            "paraphrases or noise.", top=4.7)

# ---------------------------------------------------------------- 15 thanks
s = add_slide()
t = s.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11.7), Inches(2.0))
tf = t.text_frame
p = tf.paragraphs[0]
p.text = "Thank You"
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = NAVY
q = tf.add_paragraph()
q.text = "Questions?"
q.alignment = PP_ALIGN.CENTER
q.font.size = Pt(24)

out = os.path.join(HERE, "slides.pptx")
prs.save(out)
print("wrote", out)
